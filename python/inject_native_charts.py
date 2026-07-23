#!/usr/bin/env python3
"""
inject_native_charts.py — Reads SVGs from a folder, finds data-pptx-replace-with markers,
and builds a PPTX with native charts/tables where markers exist, SVG shapes everywhere else.
"""

import sys
import json
import os
import re
from pathlib import Path
from xml.etree import ElementTree as ET

if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
if hasattr(sys.stderr, 'reconfigure'):
    sys.stderr.reconfigure(encoding='utf-8', errors='replace')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import ChartData
except ModuleNotFoundError:
    import subprocess
    print("⚠️ python-pptx missing in Python environment. Auto-installing required packages...")
    pkgs = ["python-pptx>=0.6.23", "lxml>=4.9", "Pillow>=10.0"]
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "--break-system-packages"] + pkgs)
    except Exception:
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install"] + pkgs)
        except Exception as _inst_err:
            print(f"❌ Failed to auto-install python-pptx: {_inst_err}", file=sys.stderr)
            raise

    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import ChartData
    print("✅ python-pptx auto-installed successfully.")

SLIDE_W = 13.333  # inches (16:9 presentation layout matching 1280px width)
SLIDE_H = 7.5     # inches (16:9 presentation layout matching 720px height)
PX_TO_IN = SLIDE_W / 1280  # 1px = 0.0104164 inches

def px_to_in(px):
    try:
        return float(px) * PX_TO_IN
    except Exception:
        return 0.0

def hex_to_rgb(hex_color):
    if not hex_color or not isinstance(hex_color, str):
        return RGBColor(15, 27, 56)
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c*2 for c in hex_color)
    if len(hex_color) < 6:
        return RGBColor(15, 27, 56)
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)

CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}

def add_native_chart(slide, payload):
    chart_type_key = str(payload.get("chartType", "bar")).lower()
    chart_type = CHART_TYPE_MAP.get(chart_type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)

    x = Inches(px_to_in(payload.get("x", 60)))
    y = Inches(px_to_in(payload.get("y", 150)))
    w = Inches(px_to_in(payload.get("width", 700)))
    h = Inches(px_to_in(payload.get("height", 400)))

    chart_data = ChartData()
    series_list = payload.get("series", [])

    if series_list:
        first_series = series_list[0]
        chart_data.categories = first_series.get("labels", [])
        for series in series_list:
            chart_data.add_series(
                series.get("name", "Series"),
                series.get("values", [])
            )

    chart = slide.shapes.add_chart(chart_type, x, y, w, h, chart_data).chart

    # Style the chart
    chart.has_legend = payload.get("showLegend", True)
    if chart.has_legend and hasattr(chart, 'legend') and chart.legend is not None:
        try:
            chart.legend.position = 2  # bottom
            chart.legend.include_in_layout = False
        except Exception:
            pass

    if len(chart.plots) > 0:
        plot = chart.plots[0]
        plot.has_data_labels = payload.get("showValues", False)

    # Apply colors if provided
    colors = payload.get("colors", [])
    for i, series in enumerate(chart.series):
        if i < len(colors):
            try:
                color = hex_to_rgb(colors[i])
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color
            except Exception:
                pass

def add_native_table(slide, payload):
    x = Inches(px_to_in(payload.get("x", 60)))
    y = Inches(px_to_in(payload.get("y", 150)))
    w = Inches(px_to_in(payload.get("width", 900)))
    h = Inches(px_to_in(payload.get("height", 400)))

    headers = payload.get("headers", [])
    rows = payload.get("rows", [])

    if not headers and not rows:
        return

    total_rows = len(rows) + (1 if headers else 0)
    cols = len(headers) if headers else (len(rows[0]) if rows else 1)

    table = slide.shapes.add_table(total_rows, cols, x, y, w, h).table

    header_fill = payload.get("headerFill", "0F1B38")
    header_text_color = payload.get("headerTextColor", "FFFFFF")
    row_fills = payload.get("rowFills", ["FFFFFF", "F0F4FA"])

    # Header row
    if headers:
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = str(header)
            if cell.text_frame.paragraphs:
                tf = cell.text_frame.paragraphs[0]
                if tf.runs:
                    tf.runs[0].font.bold = True
                    tf.runs[0].font.size = Pt(payload.get("fontSize", 11))
                    try:
                        tf.runs[0].font.color.rgb = hex_to_rgb(header_text_color)
                    except Exception:
                        pass
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(header_fill)
            except Exception:
                pass

    # Data rows
    for i, row in enumerate(rows):
        row_idx = i + (1 if headers else 0)
        fill_color = row_fills[i % len(row_fills)] if row_fills else "FFFFFF"
        for j, cell_val in enumerate(row):
            if j >= cols:
                break
            cell = table.cell(row_idx, j)
            cell.text = str(cell_val)
            if cell.text_frame.paragraphs:
                tf = cell.text_frame.paragraphs[0]
                if tf.runs:
                    tf.runs[0].font.size = Pt(payload.get("fontSize", 10))
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(fill_color)
            except Exception:
                pass

def safe_parse_svg_xml(svg_path):
    """Safely reads, sanitizes XML entities, auto-balances tags, and parses SVG file into ElementTree root."""
    try:
        with open(svg_path, 'r', encoding='utf-8') as f:
            content = f.read()

        content = re.sub(r'^```xml\s*', '', content)
        content = re.sub(r'^```svg\s*', '', content)
        content = re.sub(r'\s*```$', '', content)

        # Entity-aware ampersand escaping
        content = re.sub(r'&(?!amp;|lt;|gt;|quot;|apos;|#[0-9]+;|#x[0-9a-fA-F]+;)', '&amp;', content)

        ET.register_namespace("", "http://www.w3.org/2000/svg")
        try:
            return ET.fromstring(content)
        except ET.ParseError:
            # Auto-balance unclosed tags if ET parse fails initially
            open_text = len(re.findall(r'<text\b', content, re.IGNORECASE))
            close_text = len(re.findall(r'</text>', content, re.IGNORECASE))
            if open_text > close_text:
                content += "</text>" * (open_text - close_text)

            open_g = len(re.findall(r'<g\b', content, re.IGNORECASE))
            close_g = len(re.findall(r'</g>', content, re.IGNORECASE))
            if open_g > close_g:
                content += "</g>" * (open_g - close_g)

            if not re.search(r'</svg>\s*$', content, re.IGNORECASE):
                content += "\n</svg>"

            return ET.fromstring(content)
    except Exception as e:
        print(f"  ⚠️ safe_parse_svg_xml error in {os.path.basename(svg_path)}: {e}")
        return None

def parse_svg_shapes_to_pptx(svg_path, slide):
    """Fallback vector parser: parses SVG elements directly into native python-pptx shapes."""
    root = safe_parse_svg_xml(svg_path)
    if root is None:
        return 0

    def clean_tag(tag):
        return tag.split("}")[-1] if "}" in tag else tag

    shapes_added = 0

    # Process elements recursively across all nested <g> groups
    for elem in root.iter():
        tag = clean_tag(elem.tag)

        if elem.get("data-pptx-replace-with"):
            continue

        if tag == "rect":
            try:
                w = float(elem.get("width", 0))
                h = float(elem.get("height", 0))
                if w <= 0 or h <= 0:
                    continue
                x = float(elem.get("x", 0))
                y = float(elem.get("y", 0))

                fill_hex = elem.get("fill")
                stroke_hex = elem.get("stroke")
                rx = float(elem.get("rx", 0))

                shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
                shape = slide.shapes.add_shape(
                    shape_type,
                    Inches(px_to_in(x)), Inches(px_to_in(y)),
                    Inches(px_to_in(w)), Inches(px_to_in(h))
                )

                if fill_hex and fill_hex != "none":
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
                else:
                    shape.fill.background()

                if stroke_hex and stroke_hex != "none":
                    shape.line.color.rgb = hex_to_rgb(stroke_hex)
                    sw = float(elem.get("stroke-width", 1))
                    shape.line.width = Pt(sw)
                else:
                    shape.line.fill.background()

                shapes_added += 1
            except Exception:
                pass

        elif tag == "text":
            try:
                raw_text = "".join(elem.itertext()).strip()
                if not raw_text:
                    continue

                x = float(elem.get("x", 60))
                y = float(elem.get("y", 100))
                font_size = float(elem.get("font-size", 14))
                font_family = elem.get("font-family", "Calibri")
                fill_hex = elem.get("fill", "#FFFFFF")
                font_weight = elem.get("font-weight", "normal")
                text_anchor = elem.get("text-anchor", "left")

                tspans = [child for child in elem if (child.tag.split("}")[-1] if "}" in child.tag else child.tag) == "tspan"]
                num_lines = len(tspans) if tspans else max(1, len(raw_text) // 35 + 1)

                top_y = max(0.0, y - font_size * 0.85)
                box_height = font_size * 1.35 * num_lines + 12.0

                if elem.get("data-max-width"):
                    width = float(elem.get("data-max-width"))
                    left_x = x
                elif elem.get("data-width"):
                    width = float(elem.get("data-width"))
                    left_x = x
                elif text_anchor == "end":
                    width = min(1160.0, max(100.0, x - 60.0))
                    left_x = max(60.0, x - width)
                elif text_anchor in ["middle", "center"]:
                    half_avail = min(max(100.0, x - 60.0), max(100.0, 1220.0 - x))
                    width = max(100.0, min(1160.0, half_avail * 2.0))
                    left_x = max(60.0, x - width / 2.0)
                else:
                    left_x = max(60.0, min(1220.0, x))
                    width = max(100.0, min(1220.0 - left_x, 1160.0))

                txBox = slide.shapes.add_textbox(
                    Inches(px_to_in(left_x)), Inches(px_to_in(top_y)),
                    Inches(px_to_in(width)), Inches(px_to_in(box_height))
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0)
                tf.margin_right = Inches(0)
                tf.margin_top = Inches(0)
                tf.margin_bottom = Inches(0)

                if tspans:
                    for idx, ts in enumerate(tspans):
                        ts_text = "".join(ts.itertext()).strip()
                        if not ts_text:
                            continue
                        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                        p.space_before = Pt(0)
                        p.space_after = Pt(0)
                        if text_anchor == "middle":
                            p.alignment = PP_ALIGN.CENTER
                        elif text_anchor == "end":
                            p.alignment = PP_ALIGN.RIGHT
                        else:
                            p.alignment = PP_ALIGN.LEFT

                        run = p.add_run()
                        run.text = ts_text
                        run.font.name = font_family
                        run.font.size = Pt(font_size * 0.75)
                        run.font.color.rgb = hex_to_rgb(fill_hex)
                        if font_weight in ["bold", "600", "700", "800", "900"]:
                            run.font.bold = True
                else:
                    p = tf.paragraphs[0]
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    if text_anchor == "middle":
                        p.alignment = PP_ALIGN.CENTER
                    elif text_anchor == "end":
                        p.alignment = PP_ALIGN.RIGHT
                    else:
                        p.alignment = PP_ALIGN.LEFT

                    run = p.add_run()
                    run.text = raw_text
                    run.font.name = font_family
                    run.font.size = Pt(font_size * 0.75)
                    run.font.color.rgb = hex_to_rgb(fill_hex)
                    if font_weight in ["bold", "600", "700", "800", "900"]:
                        run.font.bold = True

                shapes_added += 1
            except Exception:
                pass

        elif tag == "line":
            try:
                x1 = float(elem.get("x1", 0))
                y1 = float(elem.get("y1", 0))
                x2 = float(elem.get("x2", 0))
                y2 = float(elem.get("y2", 0))
                stroke_hex = elem.get("stroke", "#F5A623")
                sw = float(elem.get("stroke-width", 2))

                w = max(abs(x2 - x1), sw)
                h = max(abs(y2 - y1), sw)
                min_x = min(x1, x2)
                min_y = min(y1, y2)

                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(px_to_in(min_x)), Inches(px_to_in(min_y)),
                    Inches(px_to_in(w)), Inches(px_to_in(h))
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(stroke_hex)
                shape.line.fill.background()

                shapes_added += 1
            except Exception:
                pass

    return shapes_added

def process_svg_file(svg_path, slide):
    """Parse one SVG, find native markers, add to slide."""
    root = safe_parse_svg_xml(svg_path)
    if root is None:
        return

    # Find all groups with data-pptx-replace-with
    for elem in root.iter():
        replace_with = elem.get("data-pptx-replace-with")
        if not replace_with:
            continue

        # Find the JSON metadata inside this group
        payload = None
        for child in elem:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "metadata":
                try:
                    text = "".join(child.itertext()).strip()
                    payload = json.loads(text)
                    break
                except json.JSONDecodeError:
                    continue

        if payload is None:
            print(f"  ⚠️ Marker '{replace_with}' in {os.path.basename(svg_path)} has no valid JSON metadata, skipping")
            continue

        print(f"  ✅ Injecting native {replace_with}: {payload.get('title', 'untitled')}")

        if replace_with == "chart":
            add_native_chart(slide, payload)
        elif replace_with == "table":
            add_native_table(slide, payload)

def add_emergency_fallback(slide, slide_num):
    """Renders a styled container and professional notice if slide content could not be rendered."""
    print(f"  🚨 Slide {slide_num}: Rendering emergency slide fallback container.")
    
    # Background container
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(SLIDE_W), Inches(SLIDE_H)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 27, 56)  # Dark executive blue
    bg.line.fill.background()

    # Content container card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(1.5),
        Inches(SLIDE_W - 2.0), Inches(4.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(26, 43, 80)
    card.line.color.rgb = RGBColor(245, 166, 35)
    card.line.width = Pt(2)

    # Title box
    tbox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(SLIDE_W - 3.0), Inches(1.0))
    p1 = tbox.text_frame.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    run1 = p1.add_run()
    run1.text = f"Slide {slide_num}: Executive Summary"
    run1.font.name = "Cambria"
    run1.font.size = Pt(24)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(255, 255, 255)

    # Body notice box
    msgbox = slide.shapes.add_textbox(Inches(1.5), Inches(3.4), Inches(SLIDE_W - 3.0), Inches(2.0))
    tf = msgbox.text_frame
    tf.word_wrap = True
    
    p2 = tf.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Content unavailable due to rendering issue.\n"
    run2.font.name = "Calibri"
    run2.font.size = Pt(16)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(245, 166, 35)

    p3 = tf.add_paragraph()
    run3 = p3.add_run()
    run3.text = "The presentation structure was generated successfully, but visual graphics for this slide could not be parsed."
    run3.font.name = "Calibri"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(128, 153, 192)

def svg_folder_to_pptx(svg_dir, output_path):
    svg_dir = Path(svg_dir)
    svg_files = sorted(svg_dir.glob("*.svg"))

    if not svg_files:
        print(f"❌ No *.svg files found in {svg_dir}")
        sys.exit(1)

    print(f"📦 Processing {len(svg_files)} slides...")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # Check for cairosvg
    cairosvg_available = False
    try:
        import cairosvg
        cairosvg_available = True
    except (ImportError, Exception):
        cairosvg_available = False

    for i, svg_file in enumerate(svg_files):
        print(f"  🎨 Slide {i+1}: {svg_file.name}")
        slide = prs.slides.add_slide(blank_layout)

        png_path = svg_file.with_suffix(".png")
        bg_added = False
        shapes_parsed = 0

        # 1. Check for pre-rendered PNG background (e.g. created by node-canvas)
        # Require at least 15KB — a valid 1280x720 slide PNG will always be much larger.
        PNG_MIN_BYTES = 15 * 1024
        if png_path.exists() and png_path.stat().st_size >= PNG_MIN_BYTES:
            try:
                pic = slide.shapes.add_picture(
                    str(png_path),
                    left=Inches(0), top=Inches(0),
                    width=Inches(SLIDE_W), height=Inches(SLIDE_H)
                )
                pic.name = f"svg_background_{i+1}"
                bg_added = True
                print(f"  🖼️ Embedded pre-rendered PNG slide background for slide {i+1}")
            except Exception as e:
                print(f"  ⚠️ Error embedding pre-rendered PNG: {e}")
        elif png_path.exists():
            print(f"  ⚠️ Skipping tiny/corrupt PNG for slide {i+1} ({png_path.stat().st_size} bytes), falling through to vector fallback")

        # 2. Fallback to CairoSVG if PNG not present
        if not bg_added and cairosvg_available:
            try:
                cairosvg.svg2png(url=str(svg_file), write_to=str(png_path), output_width=1920, output_height=1080)
                pic = slide.shapes.add_picture(
                    str(png_path),
                    left=Inches(0), top=Inches(0),
                    width=Inches(SLIDE_W), height=Inches(SLIDE_H)
                )
                pic.name = f"svg_background_{i+1}"
                bg_added = True
                if os.path.exists(png_path):
                    os.remove(png_path)
            except Exception as e:
                print(f"  ⚠️ CairoSVG background rendering error: {e}")

        # 3. Vector element fallback: parse SVG XML directly into native DrawingML shapes
        if not bg_added:
            print(f"  📐 Parsing SVG vector elements directly into PowerPoint shapes for slide {i+1}...")
            shapes_parsed = parse_svg_shapes_to_pptx(svg_file, slide)
            # Only inject native chart/table objects when in pure vector fallback mode (not when pre-rendered PNG background is used)
            process_svg_file(svg_file, slide)

        # 4. Emergency visual slide fallback if no picture or vector elements were parsed
        if not bg_added and shapes_parsed == 0:
            add_emergency_fallback(slide, i + 1)

    prs.save(output_path)
    print(f"✅ Saved: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python inject_native_charts.py <svg_dir> <output.pptx>")
        sys.exit(1)
    svg_folder_to_pptx(sys.argv[1], sys.argv[2])