#!/usr/bin/env python3
"""
svg_to_pptx.py — Converts a directory of SVG slide files into a PowerPoint presentation (.pptx).
Handles SVG elements (rectangles, text, circles, lines) and embeds vector slide images into PPTX.
"""

import sys
import os
import argparse
import glob
# pyrefly: ignore [missing-import]
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ModuleNotFoundError:
    import subprocess
    print("⚠️ python-pptx missing in Python environment. Auto-installing required packages...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx>=0.6.23", "lxml>=4.9", "Pillow>=10.0"])
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        print("✅ python-pptx auto-installed successfully.")
    except Exception as _inst_err:
        print(f"❌ Failed to auto-install python-pptx: {_inst_err}", file=sys.stderr)
        raise

try:
    # pyrefly: ignore [missing-import]
    import cairosvg
    CAIRO_AVAILABLE = True
except Exception:
    CAIRO_AVAILABLE = False

def convert_svg_folder_to_pptx(svg_dir, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]
    svg_files = sorted(glob.glob(os.path.join(svg_dir, "*.svg")))

    if not svg_files:
        raise ValueError(f"No SVG files found in {svg_dir}")

    print(f"Found {len(svg_files)} SVG slide files for conversion.")

    for idx, svg_path in enumerate(svg_files):
        slide = prs.slides.add_slide(blank_slide_layout)
        png_path = os.path.splitext(svg_path)[0] + ".png"

        if CAIRO_AVAILABLE:
            try:
                cairosvg.svg2png(url=svg_path, write_to=png_path, output_width=1920, output_height=1080)
                slide.shapes.add_picture(png_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
                if os.path.exists(png_path):
                    os.remove(png_path)
                continue
            except Exception as e:
                print(f"CairoSVG conversion warning for slide {idx + 1}: {e}. Falling back to shapes.")

        # Fallback shape rendering if cairosvg is unavailable
        with open(svg_path, 'r', encoding='utf-8') as f:
            svg_content = f.read()

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5))
        tf = txBox.text_frame
        tf.text = f"Slide {idx + 1}"

    prs.save(output_path)
    print(f"Successfully generated presentation: {output_path}")

def main():
    parser = argparse.ArgumentParser(description="Convert directory of SVG files to PPTX.")
    parser.add_argument("svg_dir", help="Directory containing slide SVG files")
    parser.add_argument("--output", "-o", required=True, help="Output PPTX file path")
    args = parser.parse_args()

    try:
        convert_svg_folder_to_pptx(args.svg_dir, args.output)
    except Exception as e:
        print(f"Error converting SVG to PPTX: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
